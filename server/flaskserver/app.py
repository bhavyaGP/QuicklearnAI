from contextlib import redirect_stdout
import unicodedata
from flask import Flask, request, jsonify, send_file
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import TextFormatter
from threading import Thread
import pyttsx3
import html
from bs4 import BeautifulSoup
from flask_cors import CORS 
import fitz
from fpdf import FPDF
import random
import re
import json
from langchain_community.llms import GPT4All  
from langchain_groq import ChatGroq
import os
from dotenv import load_dotenv
load_dotenv()
from pymongo import MongoClient
import PyPDF2
import google.generativeai as genai
import io
from google.generativeai import GenerativeModel
import jwt
from functools import wraps
from werkzeug.utils import secure_filename
import logging
from bson.objectid import ObjectId
from langchain_core.prompts import ChatPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
#from langchain_community.indexes import VectorstoreIndexCreator  # Commented out due to ImportError
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain_community.document_loaders import PyMuPDFLoader
from langchain.chains import RetrievalQA
from io import BytesIO
from PyPDF2 import PdfReader  
from langchain.schema import Document  
import chromadb
from sentence_transformers import SentenceTransformer
import google.generativeai as genai
#from langchain.document_loaders import PyPDFLoader
from pptx import Presentation
import redis
redis_client = redis.StrictRedis(host='localhost', port=6379, db=0, decode_responses=True)
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch
from agno.agent import Agent
from agno.models.groq import Groq
from agno.tools.duckduckgo import DuckDuckGoTools
import requests

app = Flask(__name__)
SECRET_KEY = "quick" 
mongo_client = MongoClient(os.getenv('MONGODB_URI', 'mongodb://localhost:27017/quicklearnai')) 
db = mongo_client["quicklearnai"]
topics_collection = db["statistics"]

CORS(app, resources={
    r"/*": {
        "origins": [
            "http://localhost:5173",  # Frontend
            "http://localhost:3000",  # Main/Proxy server
            "http://localhost:3001"   # Node server
        ],
        "methods": ["GET", "POST", "OPTIONS"],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True
    }
})

formatter = TextFormatter()

google_api_key = os.getenv("GENAI_API_KEY")
genai.configure(api_key=google_api_key)
gemini_model = genai.GenerativeModel('gemini-2.0-flash')  # Use the correct model name

# groq_api_key = "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"
groq_model = ChatGroq(
    model="llama-3.3-70b-versatile",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY")
)
os.environ["SERPER_API_KEY"] = "85a684d9cfcddab4886460954ef36f054053529b"

def get_and_enhance_transcript(youtube_url, model_type='gemini'):
    try:
        video_id = youtube_url.split('v=')[-1]
        transcript = None
        language = None

        # Fetch transcript
        for lang in ['hi', 'en']:
            try:
                transcript = YouTubeTranscriptApi().fetch(video_id, languages=[lang])
                language = lang
                break
            except:
                continue

        if not transcript:
            return None, None

        formatted_transcript = "\n".join([entry['text'] for entry in transcript])

        # Enhanced transcript prompt
        prompt = f"""
        Act as a transcript cleaner. Generate a new transcript with the same context and content as the given transcript.
        If there's a revision portion, differentiate it from the actual transcript.
        Output in sentences line by line. If the transcript lacks educational content, return 'Fake transcript'.
        Transcript: {formatted_transcript}
        """

        if model_type.lower() == 'chatgroq':
            response = groq_model.invoke(prompt)
            enhanced_transcript = response.content if hasattr(response, 'content') else str(response)
        else:  # Default to gemini
            google_api_key = os.getenv("GENAI_API_KEY")
            genai.configure(api_key=google_api_key)
            gemini_model = genai.GenerativeModel('gemini-2.0-flash') 
            response = gemini_model.generate_content(prompt)
            enhanced_transcript = response.text if hasattr(response, 'text') else str(response)
        
        return enhanced_transcript, language
    except Exception as e:
        print(f"Error in get_and_enhance_transcript: {str(e)}")
        return None, None

def generate_summary_and_quiz(transcript, num_questions, language, difficulty, model_type='gemini'):
    try:
        if 'Fake transcript' in transcript:
            return {"summary": {}, "questions": {difficulty: []}}

        prompt = f"""
        Summarize the following transcript by identifying the key topics covered, and provide a detailed summary of each topic in 6-7 sentences.
        Each topic should be labeled clearly as "Topic X", where X is the topic name. Provide the full summary for each topic in English, even if the transcript is in a different language.
        Strictly ensure that possessives (e.g., John's book) and contractions (e.g., don't) use apostrophes (') instead of quotation marks (" or "  ").

        If the transcript contains 'Fake Transcript', do not generate any quiz or summary.

        After the summary, give the name of the topic on which the transcript was all about in a maximum of 2 to 3 words.
        After summarizing, create a quiz with {num_questions} multiple-choice questions in English, based on the transcript content.
        Only generate {difficulty} difficulty questions. Format the output in JSON format as follows, just give the JSON as output, nothing before it:

        {{
            "summary": {{
                "topic1": "value1",
                "topic2": "value2",
                "topic3": "value3"
            }},
            "questions": {{
                "{difficulty}": [
                    {{
                        "question": "What is the capital of France?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Paris"
                    }},
                    {{
                        "question": "What is the capital of Germany?",
                        "options": ["Paris", "London", "Berlin", "Madrid"],
                        "answer": "Berlin"
                    }}
                ]
            }}
        }}

        Transcript: {transcript}
        """

        if model_type.lower() == 'chatgroq':
            response = groq_model.invoke(prompt)
            response_content = response.content if hasattr(response, 'content') else str(response)
        else:  # Default to gemini
            response = gemini_model.generate_content(prompt)
            response_content = response.text if hasattr(response, 'text') else str(response)

        # Extract JSON from response
        json_match = re.search(r'\{.*\}', response_content, re.DOTALL)
        if json_match:
            json_str = json_match.group(0)
            try:
                return json.loads(json_str)
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}, Raw response: {response_content}")
                return None
        else:
            print(f"No valid JSON found in response: {response_content}")
            return None
    except Exception as e:
        print(f"Error in generate_summary_and_quiz: {str(e)}")
        return None

@app.route('/quiz', methods=['POST', 'OPTIONS'])
def quiz():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    data = request.json
    youtube_link = data.get('link')
    num_questions = int(data.get('qno', 5))  # Default to 5 if not provided
    difficulty = data.get('difficulty', 'medium')  # Default to medium
    model_type = data.get('model', 'chatgroq')  # Default to gemini, can be 'chatgroq' or 'gemini'

    if not youtube_link:
        return jsonify({"error": "No YouTube URL provided"}), 400

    transcript, language = get_and_enhance_transcript(youtube_link, model_type)
    if not transcript:
        return jsonify({"error": "Failed to fetch transcript"}), 404

    summary_and_quiz = generate_summary_and_quiz(transcript, num_questions, language, difficulty, model_type)
    if summary_and_quiz:
        return jsonify(summary_and_quiz)
    else:
        return jsonify({"error": "Failed to generate quiz"}), 500

# recommendation
def validate_token_middleware():
    def middleware(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            auth_header = request.headers.get("Authorization")
            token = auth_header.split("Bearer ")[-1] if auth_header and "Bearer " in auth_header else None
            
            if not token:
                return jsonify({"message": "Unauthorized: No token provided"}), 401
            
            try:
                # Decoding the token using the correct jwt.decode()
                decoded = jwt.decode(token, SECRET_KEY, algorithms=["HS256"])
                request.user_id = decoded.get("id")
                request.user_role = decoded.get("role")  # Optional
                
                return func(*args, **kwargs)
            except jwt.ExpiredSignatureError:
                return jsonify({"message": "Unauthorized: Token has expired"}), 401
            except jwt.InvalidTokenError as e:
                print(f"Token decoding error: {e}")
                return jsonify({"message": "Unauthorized: Invalid token"}), 401
        
        return wrapper
    return middleware


from langchain.prompts import PromptTemplate
@app.route('/chat_trans', methods=['POST', 'OPTIONS'])
def chat_with_transcript():
    """Handle chat requests with YouTube transcript context"""
    if request.method == 'OPTIONS':
        return '', 204  # Handle CORS preflight request

    try:
        data = request.json
        if not data:
            return jsonify({'error': 'No JSON data provided'}), 400

        youtube_link = data.get('link')
        model_type = data.get('model', 'chatgroq')  # Default to chatgroq
        question = data.get('question')

        if not youtube_link:
            return jsonify({'error': 'Missing YouTube link'}), 400

        # Get and enhance transcript
        transcript, language = get_and_enhance_transcript(youtube_link, model_type)
        
        if "Error" in transcript:
            return jsonify({'error': transcript}), 400

        # If no question provided, just return the transcript
        if not question:
            return jsonify({
                'transcript': transcript,
                'language': language,
                'status': 'success'
            })

        # Process question with transcript context
        prompt_template = PromptTemplate(
            input_variables=["transcript", "question"],
            template="""Given the following YouTube video transcript:
            {transcript}
            
            Please answer this question based on the transcript content:
            {question}"""
        )

        formatted_prompt = prompt_template.format(
            transcript=transcript,
            question=question
        )

        # Get response from Groq
        response = groq_model.invoke(formatted_prompt)

        return jsonify({
            'answer': response.content,
            # 'transcript': transcript,
            # 'language': language,
            'status': 'success'
        })

    except Exception as e:
        return jsonify({'error': f'An error occurred: {str(e)}'}), 500

# Function to interact with LLaMA API
def llama_generate_recommendations(prompt):
    try:
        llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0,
            groq_api_key=os.getenv("GROQ_API_KEY")
        )
        
        response = llm.invoke(prompt)
        
        if hasattr(response, 'content'):
            return response.content
        else:
            return "Error: No content in response"
    except Exception as e:
        return f"Error connecting to Groq API: {e}"
 
 
 
import json

@app.route('/getonly', methods=['GET'])
@validate_token_middleware()
def get_recommendations():
    user_id = request.user_id  # Extract user ID from the token
    
    try:
        # Fetch user statistics from Redis
        statistics = redis_client.hget(f"student:{user_id}", "statistics")
        
        if not statistics:
            return jsonify({"message": "No statistics found for the provided user."}), 404
        
        # Convert JSON string to Python dictionary
        topics_data = json.loads(statistics)

        if not topics_data:
            return jsonify({"message": "No topics found for the provided user."}), 404

        # Extract only topic names
        topics_list = list(topics_data.keys())

        # Format recommendations prompt
        prompt = f"""
        Act as an intelligent recommendation generator. Based on the topics provided, generate a structured JSON response 
        with an overview, recommendations, and five YouTube video URLs for each topic. Ensure the output is in strict JSON 
        format without markdown or extra formatting. Use the following JSON structure:
        {{
            "topics": {{
                "<topic_name>": {{
                    "overview": "<brief overview>",
                    "recommendations": "<recommended steps to learn>",
                    "youtube_links": [
                        "<video_link_1>",
                        "<video_link_2>",
                        "<video_link_3>",
                        "<video_link_4>",
                        "<video_link_5>"
                    ]
                }}
            }}
        }}

        The topics are: {', '.join(topics_list)}
        """

        # Generate recommendations
        recommendations_raw = llama_generate_recommendations(prompt)

        # Ensure the response is valid JSON
        try:
            recommendations = json.loads(recommendations_raw)
        except json.JSONDecodeError:
            return jsonify({"message": "Failed to parse AI response as JSON", "raw_response": recommendations_raw}), 500

        return jsonify({
            "message": "Recommendations generated successfully",
            "recommendations": recommendations["topics"]  # Extract only relevant content
        }), 200

    except Exception as e:
        print("Error:", str(e))
        return jsonify({"message": f"An error occurred: {str(e)}"}), 500


import faiss 
from sentence_transformers import SentenceTransformer
from huggingface_hub import login
groq_api_key = os.getenv("GROQ_API_KEY")
groq_model_name = "llama3-8b-8192"
login(token=os.getenv("HUGGINGFACE_TOKEN")) 

groq_chat = ChatGroq(
    groq_api_key=groq_api_key,
    model_name=groq_model_name,
)


# Define the Groq system prompt
groq_sys_prompt = ChatPromptTemplate.from_template(
    "You are very smart at everything, you always give the best, the most accurate and most precise answers. "
    "Answer the following questions: {user_prompt}. Add more information as per your knowledge so that user can get proper knowledge, but make sure information is correct"
)
import threading
import time


embedding_model = SentenceTransformer('multi-qa-mpnet-base-cos-v1')  # Pre-trained model for embeddings
dimension = embedding_model.get_sentence_embedding_dimension()
faiss_index = faiss.IndexFlatL2(dimension) 
metadata_store = {}
pdf_storage = {}

def store_in_faiss(filename, text):
    chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
    embeddings = embedding_model.encode(chunks)
    faiss_index.add(embeddings)  
    metadata_store.update({i: filename for i in range(len(metadata_store), len(metadata_store) + len(chunks))})
genai.configure(api_key=os.getenv("GENAI_API_KEY"))
model = SentenceTransformer("all-MiniLM-L6-v2")
chroma_client = chromadb.PersistentClient(path="./chroma_db")
collection = chroma_client.get_or_create_collection(name="pdf_documents")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextToSpeechManager:
    def __init__(self):
        self.lock = threading.Lock()
    
    def speak(self, text):
        try:
            with self.lock:  # Ensure only one speech operation happens at a time
                engine = None
                try:
                    engine = pyttsx3.init()
                    engine.setProperty('rate', 150)
                    engine.setProperty('volume', 1.0)
                    engine.say(text)
                    engine.runAndWait() 
                    # print("spoke")
                    engine.startLoop(False)  # Start the event loop without blocking
                    engine.iterate()  # Process queued commands
                    engine.endLoop()  # End the event loop
                    logger.info("Speech completed successfully")
                finally:
                    if engine:
                        try:
                            engine.stop()
                        except:
                            pass
                        del engine
        except Exception as e:
            logger.error(f"Text-to-speech error: {str(e)}")
            
    def start_speaking(self, text):
        """Start a new thread for speaking"""
        thread = Thread(target=self.speak, args=(text,))
        thread.daemon = True  # Make thread daemon so it doesn't block program exit
        thread.start()
        return thread

# Create a global instance of the TTS manager
tts_manager = TextToSpeechManager()

def clean_response(text):
    """Clean and format the LLM response."""
    text = BeautifulSoup(text, "html.parser").get_text()
    text = html.unescape(text)
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = re.sub(r'[^\w\s.,!?-]', '', text)
    text = ' '.join(text.split())
    return text

def speak_text(text):
    """Convert text to speech using the TTS manager."""
    tts_manager.speak(text)

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    return " ".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
    return " ".join(text)

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file part in the request"}), 400
    
    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No file selected"}), 400
    
    file_ext = os.path.splitext(file.filename)[-1].lower()
    file_path = os.path.join("./uploads", file.filename)
    os.makedirs("./uploads", exist_ok=True)
    file.save(file_path)
    
    try:
        if file_ext == ".pdf":
            content = extract_text_from_pdf(file_path)
        elif file_ext == ".pptx":
            content = extract_text_from_pptx(file_path)
        else:
            return jsonify({"error": "Unsupported file format. Only PDF and PPTX are allowed."}), 400
        
        existing_ids = collection.get()["ids"]
        if existing_ids:
            collection.delete(ids=existing_ids)
        embedding = model.encode(content).tolist()
        collection.add(documents=[content], embeddings=[embedding], ids=[file.filename])
        
        return jsonify({"message": "File uploaded and processed successfully."}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/test-audio", methods=["GET"])
def test_audio():
    try:
        test_text = "This is a test of the text to speech system"
        logger.info("Testing text-to-speech with test message")
        
        # Start speech in a new thread
        speech_thread = tts_manager.start_speaking(test_text)
        
        return jsonify({
            "message": "Audio test initiated",
            "test_text": test_text,
            "status": "Speech initiated"
        })
    except Exception as e:
        logger.error(f"Audio test failed: {str(e)}")
        return jsonify({
            "error": "Audio test failed",
            "details": str(e)
        }), 500

@app.route("/query", methods=["POST"])
def query_file():
    try:
        data = request.get_json()
        query = data.get("query", "")
        
        logger.info(f"Received query: {query}")
        
        query_embedding = model.encode(query).tolist()
        results = collection.query(query_embeddings=[query_embedding], n_results=3)
        retrieved_texts = "\n".join(results["documents"][0])
        
        prompt = f"""
        Based on the following context, please provide a clear and concise answer to the question.
        If the answer cannot be found in the context, please say so.
        
        Context: {retrieved_texts}
        
        Question: {query}
        """
        
        response = genai.GenerativeModel("gemini-1.5-flash").generate_content(prompt)
        cleaned_response = clean_response(response.text)
        
        # Add a small delay before starting new speech
        time.sleep(0.1)  # 100ms delay
        
        speech_thread = tts_manager.start_speaking(cleaned_response)
        
        return jsonify({
            "answer": cleaned_response,
            "voice_enabled": True,
            "status": "Speech initiated"
        })
        
    except Exception as e:
        error_message = f"Error processing query: {str(e)}"
        logger.error(error_message)
        return jsonify({
            "error": error_message,
            "answer": "I apologize, but I encountered an error while processing your query. Please try again.",
            "voice_enabled": False
        }), 500


# # Configure text-to-speech settings (optional)
# @app.route("/configure-voice", methods=["POST"])
# def configure_voice():
#     try:
#         data = request.get_json()
#         rate = data.get("rate", 110)  # Default speaking rate
#         volume = data.get("volume", 1.0)  # Default volume
#         voice_id = data.get("voice_id")  # Voice identifier
        
#         engine.setProperty('rate', rate)
#         engine.setProperty('volume', volume)
        
#         if voice_id:
#             voices = engine.getProperty('voices')
#             for voice in voices:
#                 if voice.id == voice_id:
#                     engine.setProperty('voice', voice.id)
#                     break
        
#         return jsonify({"message": "Voice settings updated successfully"}), 200
#     except Exception as e:
#         return jsonify({"error": f"Error configuring voice: {str(e)}"}), 500

# Configure text-to-speech settings (optional)
# @app.route("/configure-voice", methods=["POST"])
# def configure_voice():
#     try:
#         data = request.get_json()
#         rate = data.get("rate", 110)  # Default speaking rate
#         volume = data.get("volume", 1.0)  # Default volume
#         voice_id = data.get("voice_id")  # Voice identifier
        
#         engine.setProperty('rate', rate)
#         engine.setProperty('volume', volume)
        
#         if voice_id:
#             voices = engine.getProperty('voices')
#             for voice in voices:
#                 if voice.id == voice_id:
#                     engine.setProperty('voice', voice.id)
#                     break
        
#         return jsonify({"message": "Voice settings updated successfully"}), 200
#     except Exception as e:
#         return jsonify({"error": f"Error configuring voice: {str(e)}"}), 500
# # MindMap

def fetch_youtube_transcript(video_url):
    try:
        video_id = video_url.split("v=")[-1]
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'hi'])
        return " ".join([entry["text"] for entry in transcript])  # Clean transcript
    except Exception as e:
        return {"error": f"Error fetching transcript: {str(e)}"}

def generate_mind_map(content):
    prompt = f"""
    Extract key concepts from the following text and structure them into a JSON-based mind map.
    Organize it into: "Topic" -> "Subtopics" -> "Details".

    Text: {content}

    Output **ONLY** valid JSON in this format (no extra text, no explanations):
    {{
        "topic": "Main Topic",
        "subtopics": [
            {{"name": "Subtopic 1", "details": ["Detail 1", "Detail 2"]}},
            {{"name": "Subtopic 2", "details": ["Detail 3", "Detail 4"]}}
        ]
    }}
    """

    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0,
        groq_api_key=os.getenv("GROQ_API_KEY"),
    )      

    response = llm.invoke(prompt)

    # Ensure response is a string
    raw_json = response.content.strip() if hasattr(response, "content") else str(response)

    # Remove unwanted formatting (like triple backticks and newlines)
    cleaned_json_str = raw_json.replace("```json", "").replace("```", "").replace("\n", "").strip()

    # Convert to valid JSON
    try:
        return json.loads(cleaned_json_str)
    except json.JSONDecodeError:
        return {"error": f"Invalid JSON response: {cleaned_json_str}"}

@app.route("/generate_mind_map", methods=['GET'])
def generate_mind_map_endpoint():
    # print("✅ Endpoint called!")  # Debugging
    video_url = request.args.get('video_url')

    if not video_url:
        return jsonify({"error": "No video URL provided"}), 400

    transcript = fetch_youtube_transcript(video_url)
    if isinstance(transcript, dict) and "error" in transcript:
        return jsonify(transcript), 400

    mind_map = generate_mind_map(transcript)
   
    return jsonify(mind_map)

llm = ChatGroq(
    model="llama-3.3-70b-versatile",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY"),
)

def generate_quiz(topic: str, num_questions: int, difficulty: str):
    """Generate a quiz based on the given topic."""
    prompt = f"""
    Create a quiz on the topic: "{topic}". Generate {num_questions} multiple-choice questions.
    The questions should be of {difficulty} difficulty.
    Format the output strictly in JSON format as follows:
    
    {{
       
        "questions": {{
            "{difficulty}": [
                {{
                    "question": "What is ...?",
                    "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
                    "answer": "Option 1"
                }}
            ]
        }}
    }}
    """
    response = llm.invoke(prompt)
    return response.content if hasattr(response, 'content') else response.text

@app.route("/llm_quiz", methods=["POST"])
def quiz_endpoint():
    data = request.json
    topic = data.get("topic")
    num_questions = data.get("num_questions")
    difficulty = data.get("difficulty")
    
    if not topic:
        return jsonify({"error": "Topic is required"}), 400
    
    try:
        response_content = generate_quiz(topic, num_questions, difficulty)

        
        try:
            result = json.loads(response_content)
        except json.JSONDecodeError:
            json_start = response_content.find('{')
            json_end = response_content.rfind('}') + 1
            if json_start != -1 and json_end != -1:
                json_str = response_content[json_start:json_end]
                result = json.loads(json_str)
            else:
                return jsonify({"error": "Could not parse JSON from response"}), 500
        
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500




llm = ChatGroq(
    model="llama-3.3-70b-versatile",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY")
)

def generate_quiz(topic: str, num_questions: int, difficulty: str):
    """Generate a quiz based on the given topic."""
    prompt = f"""
    Create a quiz on the topic: "{topic}". Generate {num_questions} multiple-choice questions.
    The questions should be of {difficulty} difficulty.
    Format the output strictly in JSON format as follows:
    
    {{
       
        "questions": {{
            "{difficulty}": [
                {{
                    "question": "What is ...?",
                    "options": ["Option 1", "Option 2", "Option 3", "Option 4"],
                    "answer": "Option 1"
                }}
            ]
        }}
    }}
    """
    response = llm.invoke(prompt)
    return response.content if hasattr(response, 'content') else response.text



UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "generated_papers"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

groq_api_key = os.getenv("GROQ_API_KEY")

def extract_questions_from_pdf(pdf_path):
    """
    Extract questions from a PDF file more robustly.
    Handles various question formats and improves text extraction.
    """
    try:
        doc = fitz.open(pdf_path)
        questions = []
        current_question = ""
        
        for page in doc:
            # Extract text with different methods for robustness
            text = page.get_text("text")  # Primary method
            if not text.strip():  # Fallback if text extraction fails
                text = page.get_text("blocks")  # Extract by blocks if needed
                text = "\n".join(block[4] for block in text if len(block) > 4)  # block[4] is the text content
            
            # Split text into lines and clean up
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            for line in lines:
                # Check if line could be the start of a new question
                is_new_question = False
                # Pattern 1: Starts with number followed by dot or parenthesis (e.g., "1.", "1)")
                if re.match(r'^\d+[.\)]\s', line):
                    is_new_question = True
                # Pattern 2: Contains a question mark and is long enough to be meaningful
                elif '?' in line and len(line) > 10:
                    is_new_question = True
                
                if is_new_question:
                    # Save previous question if exists
                    if current_question:
                        questions.append(current_question.strip())
                    current_question = line
                elif current_question:  # Append to existing question (multi-line)
                    current_question += " " + line
            
            # Append the last question on the page if exists
            if current_question:
                questions.append(current_question.strip())
                current_question = ""  # Reset for next page
        
        # Final filtering of questions
        filtered_questions = []
        for q in questions:
            # Ensure question is meaningful: has a question mark and sufficient length
            if '?' in q and len(q) > 10:
                # Clean up extra spaces and ensure proper encoding
                cleaned_q = " ".join(q.split())
                filtered_questions.append(cleaned_q)
        
        # If no questions found, try a more aggressive extraction as last resort
        if not filtered_questions:
            all_text = ""
            for page in doc:
                all_text += page.get_text("text") + "\n"
            # Split by question marks and filter
            potential_questions = [q.strip() for q in all_text.split('?') if q.strip()]
            for pq in potential_questions:
                if len(pq) > 10:
                    cleaned_q = f"{pq}?".strip()  # Re-add question mark
                    filtered_questions.append(cleaned_q)
        
        doc.close()
        return filtered_questions if filtered_questions else ["No valid questions found in PDF"]

    except Exception as e:
        return [f"Error extracting questions: {str(e)}"]

def generate_questions(extracted_questions, num_questions):
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        temperature=0.7,
        groq_api_key=groq_api_key
    )
    
    prompt = (
        "Based on the following sample questions, generate similar academic questions that are clear, meaningful, "
        "and properly formatted. Each question should end with a question mark and be suitable for an exam paper:\n\n"
        f"Sample questions:\n{chr(10).join(extracted_questions[:5])}\n\n"
        f"Generate {num_questions} new questions in a similar style and complexity level."
    )
    
    response = llm.invoke(prompt)
    generated_text = response.content
    
    new_questions = [q.strip() for q in generated_text.split('\n') if q.strip()]
    valid_questions = []
    for q in new_questions:
        if q and q[0].isdigit() and '?' in q and len(q) > 15:
            valid_questions.append(q)
        elif q and '?' in q and len(q) > 15:
            valid_questions.append(f"{len(valid_questions) + 1}. {q}")
    
    return valid_questions[:num_questions]

def create_question_paper(questions, filename, set_number):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    
    # Set light blue background
    pdf.set_fill_color(240, 248, 255)  # Light blue RGB
    pdf.rect(0, 0, pdf.w, pdf.h, 'F')
    
    # Add heading
    pdf.set_font("Arial", "B", 16)
    pdf.set_text_color(0, 0, 0)  # Black text
    pdf.cell(0, 10, f"Set Number - {set_number}", ln=True, align='C')
    pdf.ln(10)
    
    # Add questions with proper numbering
    pdf.set_font("Arial", size=12)
    for idx, question in enumerate(questions, 1):
        q_text = question.split('.', 1)[1].strip() if '.' in question[:3] else question
        safe_question = f"{idx}. {q_text}".encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, safe_question)
        pdf.ln(5)
    
    # Add footer text at the bottom
    pdf.set_y(-20)  # Position 20mm from bottom
    pdf.set_font("Arial", "I", 8)
    pdf.set_text_color(100, 100, 100)  # Gray text
    pdf.cell(0, 10, "Powered by QuickLearn AI", ln=True, align='C')
    
    pdf.output(filename)

@app.route('/paper_upload', methods=['POST'])
def upload_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files['file']
    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)
    return jsonify({"message": "File uploaded successfully", "file_path": file_path})

@app.route('/generate_paper', methods=['POST'])
def generate_papers():
    data = request.json
    pdf_path = data.get("file_path")
    num_questions = data.get("num_questions", 10)
    num_papers = data.get("num_papers", 1)
    
    if not pdf_path or not os.path.exists(pdf_path):
        return jsonify({"error": "Invalid file path"}), 400
    
    try:
        extracted_questions = extract_questions_from_pdf(pdf_path)
        if not extracted_questions:
            return jsonify({"error": "No valid questions found in PDF"}), 400
        
        generated_questions = generate_questions(extracted_questions, num_questions * num_papers)
        all_questions = extracted_questions + generated_questions
        random.shuffle(all_questions)
        
        pdf_paths = []
        for i in range(num_papers):
            start_idx = i * num_questions
            end_idx = start_idx + num_questions
            if start_idx >= len(all_questions):
                break
                
            selected_questions = all_questions[start_idx:end_idx]
            if len(selected_questions) < num_questions:
                remaining = num_questions - len(selected_questions)
                extra_questions = generate_questions(extracted_questions, remaining)
                selected_questions.extend(extra_questions)
            
            paper_path = os.path.join(OUTPUT_FOLDER, f"question_paper_set_{i+1}.pdf")
            create_question_paper(selected_questions, paper_path, i + 1)
            pdf_paths.append(paper_path)
        
        return jsonify({"message": "Papers generated", "files": pdf_paths})
    
    except Exception as e:
        return jsonify({"error": f"An error occurred: {str(e)}"}), 500

@app.route('/download/<filename>', methods=['GET'])
def download_paper(filename):
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    return jsonify({"error": "File not found"}), 404

@app.route('/', methods=['GET'])
def health():
    return jsonify({"status": "ok"}) 

def initialize_question_bank_agent():
    try:
        groq_model = Groq(
            id="llama3-70b-8192",
            api_key=groq_api_key
        )
        agent = Agent(
            model=groq_model,
            description=(
                "You are a helpful assistant that creates challenging practice problems. "
                "When asked to create questions, format them with numbers (1., 2., etc.) "
                "and make them clear and well-structured."
            ),
            tools=[DuckDuckGoTools()],
            show_tool_calls=False,
            markdown=False
        )
        return agent
    except Exception as e:
        raise Exception(f"Error initializing agent: {str(e)}")

def get_agent_response(agent, prompt):
    try:
        buffer = io.StringIO()
        with redirect_stdout(buffer):
            agent.print_response(prompt)
        response = buffer.getvalue()
        
        if not response:
            raise Exception("No response received from agent")
        
        cleaned = unicodedata.normalize('NFKD', response).encode('ascii', 'ignore').decode('ascii')
        cleaned = re.sub(r'\x1B(?:[@-Z\\-_]|\[[0-?]*[ -/]*[@-~])', '', cleaned)
        cleaned = re.sub(r'Running:.*?\n', '', cleaned)
        cleaned = re.sub(r'Response:?\s*', '', cleaned)
        cleaned = re.sub(r'<tool-use>.*?</tool-use>', '', cleaned, flags=re.DOTALL)
        cleaned = re.sub(r'\(\d+\.\d+s\)', '', cleaned)
        cleaned = re.sub(r'Here are \d+ practice problems.*?:(?=\n)', '', cleaned)
        cleaned = re.sub(r'1\., 2\., etc\.\).*?\n', '', cleaned)
        cleaned = re.sub(r'Make each question clear and well-formatted.*?numerical questions over theory\.', '', cleaned)
        cleaned = re.sub(r'\n\s*\n', '\n\n', cleaned)
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = cleaned.strip()
        
        questions = []
        current_question = ""
        for line in cleaned.split('\n'):
            line = line.strip()
            if not line:
                continue
            if re.match(r'^\d+\.', line):
                if current_question:
                    questions.append(current_question.strip())
                current_question = line
            elif current_question:
                current_question += " " + line
        
        if current_question:
            questions.append(current_question.strip())
        
        if not questions:
            matches = re.findall(r'(\d+\..*?)(?=\d+\.|$)', cleaned, re.DOTALL)
            questions = [q.strip() for q in matches]
        
        if not questions:
            return "1. Error: No questions could be generated. Please try again."
        
        formatted_questions = []
        for i, q in enumerate(questions, 1):
            q = q.strip()
            if not re.match(r'^\d+\.', q):
                q = f"{i}. {q}"
            formatted_questions.append(q)
        
        return "\n\n".join(formatted_questions)
        
    except Exception as e:
        return f"1. An error occurred while generating questions: {str(e)}"

def create_question_bank_pdf(text, subject):
    question_bank_dir = os.path.join(os.getcwd(), "generated_papers")
    os.makedirs(question_bank_dir, exist_ok=True)
    
    filename = os.path.join(question_bank_dir, f"{subject.replace(' ', '_')}_Questions.pdf")
    
    try:
        doc = SimpleDocTemplate(filename, pagesize=letter,
                              rightMargin=72, leftMargin=72,
                              topMargin=72, bottomMargin=72)
        
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=28,
            spaceAfter=20,
            alignment=1,
            textColor=HexColor('#FFD700'),
            leading=32
        )
        
        question_style = ParagraphStyle(
            'CustomQuestion',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=12,
            leading=18,
            spaceBefore=20,
            spaceAfter=20,
            firstLineIndent=0,
            leftIndent=20,
            textColor=HexColor('#00FF9D')
        )
        
        elements = []
        elements.append(Paragraph(f"{subject} Practice Questions", title_style))
        elements.append(Spacer(1, 0.2*inch))
        
        questions = [q.strip() for q in text.split('\n\n') if q.strip()]
        for q in questions:
            q = q.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            p = Paragraph(q, question_style)
            elements.append(p)
        
        def add_watermark_and_background(canvas, doc):
            canvas.saveState()
            canvas.setFillColorRGB(0, 0, 0)
            canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=1)
            canvas.setFont("Helvetica-Bold", 60)
            canvas.setFillColorRGB(0.9, 0.9, 0.9, alpha=0.5)
            canvas.translate(doc.pagesize[0]/2, doc.pagesize[1]/2)
            canvas.rotate(45)
            canvas.drawCentredString(0, 0, "QuickLearn AI")
            canvas.restoreState()
        
        doc.build(elements, onFirstPage=add_watermark_and_background, onLaterPages=add_watermark_and_background)
        
        return filename
    
    except Exception as e:
        raise Exception(f"Error building PDF: {str(e)}")

# Add the new route
@app.route('/question_bank', methods=['POST', 'OPTIONS'])
def generate_question_bank():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200
        
    try:
        if not request.is_json:
            return jsonify({"error": "Request must be JSON"}), 400
        
        data = request.get_json()
        subject = data.get('topic')
        
        if not subject or not isinstance(subject, str):
            return jsonify({"error": "Invalid or missing 'topic' in payload"}), 400
        
        subject = subject.strip()
        if not subject:
            return jsonify({"error": "Topic cannot be empty"}), 400
        
        agent = initialize_question_bank_agent()
        
        prompt = (
            f"Create 10 challenging practice problems on {subject}. "
            "Number each question with a number and period (1., 2., etc.). "
            "Make each question clear and well-formatted on its own line. "
            "Include a mix of difficulty levels from basic to advanced. "
            "Focus on numerical questions over theoretical ones."
        )
 
        result_text = get_agent_response(agent, prompt)
        
        if result_text.startswith("1. Error:") or result_text.startswith("1. An error occurred"):
            return jsonify({"error": "Failed to generate questions"}), 500
        
        pdf_path = create_question_bank_pdf(result_text, subject)
        
        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=f"{subject.replace(' ', '_')}_Questions.pdf",
            mimetype='application/pdf'
        )
    
    except Exception as e:
        return jsonify({"error": f"Server error: {str(e)}"}), 500
    
def search_youtube_videos(topic, max_results=3):
    url = "https://google.serper.dev/videos"
    payload = {"q": f"{topic} tutorial"}
    headers = {
        "X-API-KEY": os.environ["SERPER_API_KEY"],
        "Content-Type": "application/json"
    }
    try:
        response = requests.post(url, json=payload, headers=headers)
        response.raise_for_status()
        data = response.json()
        
        if "videos" not in data:
            return []
        
        # Extract up to max_results URLs
        urls = [video.get("link", "") for video in data.get("videos", [])[:max_results]]
        return urls
    except requests.RequestException as e:
        print(f"Serper API error for {topic}: {e}")
        return []

def is_valid_youtube_url(url):
    pattern = r'^(https?://(www\.)?youtube\.com/watch\?v=[\w-]{11}|https?://youtu\.be/[\w-]{11})'
    return bool(re.match(pattern, url))

@app.route('/youtube_videos', methods=['POST', 'OPTIONS'])
def youtube_videos():
    if request.method == 'OPTIONS':
        return jsonify({'status': 'ok'}), 200

    try:
        data = request.json
        if not data or 'topic' not in data:
            return jsonify({
                "success": False,
                "error": "Missing 'topics' in JSON body"
            }), 400
        
        topics = data['topic']
        
        # Convert single topic to list if needed
        if isinstance(topics, str):
            topics = [topics]
        
        if not isinstance(topics, list) or not topics:
            return jsonify({
                "success": False,
                "error": "'topics' must be a non-empty list"
            }), 400
        
        result = {}
        for topic in topics:
            video_urls = search_youtube_videos(topic, max_results=3)
            valid_urls = [url for url in video_urls if is_valid_youtube_url(url)]
            result[topic] = valid_urls[:3]
        
        return jsonify({
            "success": True,
            "data": result
        })
    
    except ValueError as e:
        return jsonify({
            "success": False,
            "error": "Invalid JSON format"
        }), 400
    except Exception as e:
        print(f"Server error: {e}")
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500

if __name__ == '__main__':
    app.run(debug=True, port=5001)
    